"""
Generate PowerPoint presentation for Depression Text Analysis project.
Uses dark text colors for visibility on white/light backgrounds.
"""
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx library not installed.")
    print("Install it using: pip install python-pptx")
    sys.exit(1)

# Create presentation with widescreen aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Modern Color Palette - Vibrant and Premium
PRIMARY_GRADIENT_START = RGBColor(88, 86, 214)    # Deep purple
PRIMARY_GRADIENT_END = RGBColor(52, 138, 199)      # Ocean blue
TITLE_COLOR = RGBColor(255, 255, 255)              # White for contrast on gradients
ACCENT_CORAL = RGBColor(255, 107, 107)             # Vibrant coral
ACCENT_TEAL = RGBColor(72, 219, 251)               # Bright teal  
ACCENT_GOLD = RGBColor(255, 195, 113)              # Warm gold
TEXT_COLOR = RGBColor(45, 55, 72)                  # Deep charcoal
SUBTITLE_COLOR = RGBColor(113, 128, 150)           # Cool gray
LIGHT_BG = RGBColor(247, 250, 252)                 # Very light blue-gray
CARD_BG = RGBColor(255, 255, 255)                  # White for cards

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Gradient background
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45
    fill.gradient_stops[0].color.rgb = PRIMARY_GRADIENT_START
    fill.gradient_stops[1].color.rgb = PRIMARY_GRADIENT_END
    
    # Decorative accent shapes
    # Top right accent circle
    accent1 = slide.shapes.add_shape(
        1, Inches(10.5), Inches(-1), Inches(4), Inches(4)
    )
    accent1.fill.solid()
    accent1.fill.fore_color.rgb = ACCENT_CORAL
    accent1.line.fill.background()
    accent1.fill.transparency = 0.3
    
    # Bottom left accent
    accent2 = slide.shapes.add_shape(
        1, Inches(-1), Inches(5.5), Inches(3.5), Inches(3.5)
    )
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = ACCENT_TEAL
    accent2.line.fill.background()
    accent2.fill.transparency = 0.4
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.3), Inches(1.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Decorative line under title
    line = slide.shapes.add_shape(
        1, Inches(5), Inches(4.3), Inches(3.3), Inches(0.15)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_GOLD
    line.line.fill.background()
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.6), Inches(10.3), Inches(1.5))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(26)
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True
    
    return slide

def add_content_slide(prs, title, bullets, image_path=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Light gradient background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Title bar with gradient
    title_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    title_bar_fill = title_bar.fill
    title_bar_fill.gradient()
    title_bar_fill.gradient_angle = 0
    title_bar_fill.gradient_stops[0].color.rgb = PRIMARY_GRADIENT_START
    title_bar_fill.gradient_stops[1].color.rgb = PRIMARY_GRADIENT_END
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Determine layout based on image
    if image_path and os.path.exists(image_path):
        content_width = Inches(6.2)
        content_left = Inches(0.5)
        
        # Add image with subtle border
        img_left = Inches(7.2)
        img_top = Inches(1.5)
        img_width = Inches(5.6)
        
        # Image background card
        img_card = slide.shapes.add_shape(
            1, img_left - Inches(0.15), img_top - Inches(0.15), 
            img_width + Inches(0.3), Inches(5.3)
        )
        img_card.fill.solid()
        img_card.fill.fore_color.rgb = CARD_BG
        img_card.line.color.rgb = RGBColor(226, 232, 240)
        img_card.line.width = Pt(2)
        img_card.shadow.inherit = False
        
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
    else:
        content_width = Inches(12.3)
        content_left = Inches(0.5)
        if image_path:
            missing_images.append(image_path)
    
    # Content card
    card = slide.shapes.add_shape(
        1, content_left, Inches(1.5), content_width, Inches(5.3)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = RGBColor(226, 232, 240)
    card.line.width = Pt(2)
    card.shadow.inherit = False
    
    # Bullets
    if bullets:
        text_box = slide.shapes.add_textbox(
            content_left + Inches(0.4), Inches(1.8), 
            content_width - Inches(0.8), Inches(4.7)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # Style bullets with colored markers
            if bullet.strip().startswith('•'):
                p.text = bullet
                p.font.size = Pt(20)
                p.font.color.rgb = TEXT_COLOR
                p.space_after = Pt(10)
                p.level = 0
            elif bullet.strip() == "":
                p.text = ""
                p.space_after = Pt(6)
            elif bullet.strip().startswith('STEP') or bullet.strip().startswith('KEY') or bullet.strip().startswith('SUMMARY'):
                # Section headers
                p.text = bullet
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_GRADIENT_START
                p.space_after = Pt(8)
            else:
                # Sub-bullets or regular text
                p.text = bullet
                p.font.size = Pt(19)
                p.font.color.rgb = TEXT_COLOR
                p.space_after = Pt(8)
    
    return slide

def add_image_slide(prs, title, explanation, image_path):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Light background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Title bar with gradient
    title_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(1.0)
    )
    title_bar_fill = title_bar.fill
    title_bar_fill.gradient()
    title_bar_fill.gradient_angle = 0
    title_bar_fill.gradient_stops[0].color.rgb = PRIMARY_GRADIENT_START
    title_bar_fill.gradient_stops[1].color.rgb = PRIMARY_GRADIENT_END
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.15), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Explanation with accent background
    exp_bg = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.55)
    )
    exp_bg.fill.solid()
    exp_bg.fill.fore_color.rgb = ACCENT_TEAL
    exp_bg.fill.transparency = 0.15
    exp_bg.line.fill.background()
    
    exp_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.3), Inches(0.5))
    tf = exp_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = explanation
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_COLOR
    p.font.italic = True
    
    # Image with frame
    if os.path.exists(image_path):
        # Image card
        img_card = slide.shapes.add_shape(
            1, Inches(0.7), Inches(1.9), Inches(11.9), Inches(5.2)
        )
        img_card.fill.solid()
        img_card.fill.fore_color.rgb = CARD_BG
        img_card.line.color.rgb = RGBColor(226, 232, 240)
        img_card.line.width = Pt(3)
        
        slide.shapes.add_picture(image_path, Inches(0.85), Inches(2.05), width=Inches(11.6))
    else:
        missing_images.append(image_path)
    
    return slide

def add_two_image_slide(prs, title, explanation, img1_path, img2_path, label1, label2):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Light background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Title bar with gradient
    title_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(1.0)
    )
    title_bar_fill = title_bar.fill
    title_bar_fill.gradient()
    title_bar_fill.gradient_angle = 0
    title_bar_fill.gradient_stops[0].color.rgb = PRIMARY_GRADIENT_START
    title_bar_fill.gradient_stops[1].color.rgb = PRIMARY_GRADIENT_END
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.15), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Explanation with background
    exp_bg = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.5)
    )
    exp_bg.fill.solid()
    exp_bg.fill.fore_color.rgb = ACCENT_CORAL
    exp_bg.fill.transparency = 0.15
    exp_bg.line.fill.background()
    
    exp_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.3), Inches(0.45))
    tf = exp_box.text_frame
    p = tf.paragraphs[0]
    p.text = explanation
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_COLOR
    p.font.italic = True
    
    # Image 1 with card
    if os.path.exists(img1_path):
        # Card background
        card1 = slide.shapes.add_shape(
            1, Inches(0.4), Inches(1.85), Inches(6.2), Inches(5.0)
        )
        card1.fill.solid()
        card1.fill.fore_color.rgb = CARD_BG
        card1.line.color.rgb = ACCENT_TEAL
        card1.line.width = Pt(3)
        
        slide.shapes.add_picture(img1_path, Inches(0.5), Inches(1.95), width=Inches(6.0))
        
        # Label with gradient background
        lbl_bg1 = slide.shapes.add_shape(
            1, Inches(0.4), Inches(6.5), Inches(6.2), Inches(0.5)
        )
        lbl_bg1.fill.solid()
        lbl_bg1.fill.fore_color.rgb = ACCENT_TEAL
        lbl_bg1.line.fill.background()
        
        lbl1 = slide.shapes.add_textbox(Inches(0.4), Inches(6.55), Inches(6.2), Inches(0.4))
        tf = lbl1.text_frame
        p = tf.paragraphs[0]
        p.text = label1
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
    else:
        missing_images.append(img1_path)
    
    # Image 2 with card
    if os.path.exists(img2_path):
        # Card background
        card2 = slide.shapes.add_shape(
            1, Inches(6.7), Inches(1.85), Inches(6.2), Inches(5.0)
        )
        card2.fill.solid()
        card2.fill.fore_color.rgb = CARD_BG
        card2.line.color.rgb = ACCENT_CORAL
        card2.line.width = Pt(3)
        
        slide.shapes.add_picture(img2_path, Inches(6.8), Inches(1.95), width=Inches(6.0))
        
        # Label with gradient background
        lbl_bg2 = slide.shapes.add_shape(
            1, Inches(6.7), Inches(6.5), Inches(6.2), Inches(0.5)
        )
        lbl_bg2.fill.solid()
        lbl_bg2.fill.fore_color.rgb = ACCENT_CORAL
        lbl_bg2.line.fill.background()
        
        lbl2 = slide.shapes.add_textbox(Inches(6.7), Inches(6.55), Inches(6.2), Inches(0.4))
        tf = lbl2.text_frame
        p = tf.paragraphs[0]
        p.text = label2
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
    else:
        missing_images.append(img2_path)
    
    return slide


# ============== BUILD PRESENTATION ==============

# Get the script directory to ensure correct paths
script_dir = os.path.dirname(os.path.abspath(__file__))
img_dir = os.path.join(script_dir, "output", "figures")

# Track missing images for reporting at the end
missing_images = []

# Slide 1: Title
add_title_slide(prs, 
    "Depression Text Analysis",
    "Word Frequency & Correlation Study\nAnalyzing Linguistic Patterns in Clinical Interviews")

# Slide 2: Project Overview
add_content_slide(prs, "Project Overview", [
    "• GOAL: Identify linguistic patterns that correlate with depression",
    "",
    "• Analyze psychologist-patient interview transcripts from clinical settings",
    "• Apply Natural Language Processing (NLP) techniques to extract patterns",
    "• Find specific words and phrases that distinguish depressed from non-depressed individuals",
    "",
    "• RESEARCH QUESTION:",
    "  'Can specific words and phrases serve as indicators of depression?'",
    "",
    "• APPROACH: Progressive complexity",
    "  - Start simple: Basic word frequency counts",
    "  - Progress to advanced: Sentiment analysis, statistical testing, machine learning"
])

# Slide 3: Dataset Description
add_content_slide(prs, "Dataset: AVEC 2017 Depression Challenge", [
    "• 189 participants interviewed by virtual agent 'Ellie'",
    "• Semi-structured clinical interviews averaging 20 minutes",
    "• Participants answered questions about feelings, relationships, activities",
    "",
    "CLASSIFICATION:",
    "• 57 classified as DEPRESSED (PHQ-8 score ≥ 10)",
    "• 132 classified as NON-DEPRESSED (PHQ-8 < 10)",
    "",
    "• PHQ-8: Patient Health Questionnaire (8 items)",
    "  - Standard depression screening tool",
    "  - Scores range 0-24 (higher = more severe symptoms)",
    "  - Threshold of 10+ indicates clinical depression",
    "",
    "• Average transcript length: ~756 words per participant"
])

# Slide 4: PHQ Distribution Figure
add_image_slide(prs, "PHQ-8 Score Distribution",
    "Distribution shows clear separation at threshold 10. Most participants cluster at low scores (healthy) or moderate-high scores (depressed), with fewer in the middle range.",
    f"{img_dir}/04_phq_distribution.png")

# Slide 5: Dataset Characteristics
add_content_slide(prs, "Dataset Characteristics & Challenges", [
    "KEY STATISTICS:",
    "• Total participants: 189",
    "• Class imbalance: 30% depressed vs 70% non-depressed",
    "• Vocabulary size: 8,474 unique words after preprocessing",
    "• Average words per transcript: 756 words",
    "",
    "CHALLENGES:",
    "• Imbalanced classes require careful statistical analysis",
    "• Conversational speech includes many fillers (um, uh)",
    "• Individual variation in speaking style",
    "• Need to distinguish depression markers from general interview patterns"
])

# Slide 6: Methodology Overview
add_content_slide(prs, "Analysis Pipeline", [
    "STEP 1: Load Data",
    "• Read interview transcripts and PHQ-8 labels",
    "• Binary classification: Depressed (1) vs Non-depressed (0)",
    "",
    "STEP 2: Text Preprocessing",
    "• Lowercase conversion, remove punctuation",
    "• Remove stopwords and interview fillers",
    "• Tokenization into individual words",
    "",
    "STEP 3-4: Frequency Analysis",
    "• Count word occurrences across all transcripts",
    "• Compare frequencies between groups",
    "",
    "STEP 5-9: Advanced Analysis",
    "• Statistical correlation testing (point-biserial)",
    "• TF-IDF weighting • N-gram extraction • Sentiment analysis"
])

# Slide 7: Preprocessing Details
add_content_slide(prs, "Text Preprocessing Methodology", [
    "CLEANING STEPS:",
    "• Convert all text to lowercase for normalization",
    "• Remove punctuation and special characters",
    "• Remove common English stopwords (the, a, to, is, etc.)",
    "• Remove interview fillers specific to this dataset:",
    "  - Fillers: um, uh, yeah, like, you know",
    "  - These add no semantic meaning for depression detection",
    "• Tokenize: Split text into individual words",
    "",
    "EXAMPLE TRANSFORMATION:",
    "  Original: 'I'm really feeling um like I don't know...'",
    "  Processed: ['im', 'really', 'feeling', 'dont', 'know']",
    "",
    "RESULT: Reduced from raw text to 8,474 unique meaningful tokens"
])

# Slide 8: Word Frequency Analysis - Method
add_content_slide(prs, "Method: Word Frequency Analysis", [
    "APPROACH:",
    "• Count occurrence of each word across all 189 transcripts",
    "• Use Counter from Python collections library",
    "• Identify the most common words in the dataset",
    "",
    "PURPOSE:",
    "• Understand overall vocabulary and speech patterns",
    "• Identify baseline of common conversational words",
    "• Establish context before group comparisons",
    "",
    "EXPECTED FINDINGS:",
    "• Most frequent words likely to be general conversational terms",
    "• Casual speech markers and common verbs/adjectives",
    "• These serve as baseline for comparison with group-specific patterns"
])

# Slide 9: Word Frequency - Results
add_image_slide(prs, "Results: Most Common Words Overall",
    "Bar chart shows the 30 most frequent words across all interviews. Common conversational words dominate: 'just', 'dont', 'know', 'really'. This establishes our baseline vocabulary before examining group differences.",
    f"{img_dir}/01_top_words.png")

# Slide 10: Group Comparison - Method
add_content_slide(prs, "Method: Group Frequency Comparison", [
    "APPROACH:",
    "• Separate transcripts by depression status (Depressed vs Non-Depressed)",
    "• Calculate average word frequency for each group",
    "• Identify top 30 words for each group independently",
    "• Visualize side-by-side for direct comparison",
    "",
    "STATISTICAL CONSIDERATION:",
    "• Account for class imbalance (57 vs 132 participants)",
    "• Use normalized frequencies (average per person)",
    "• Focus on relative differences, not absolute counts",
    "",
    "HYPOTHESIS:",
    "• Depressed group may use more negative emotion words",
    "• Depressed group may reference mental health treatment",
    "• Non-depressed may use more active/social words"
])

# Slide 11: Group Comparison - Results
add_image_slide(prs, "Results: Depressed vs Non-Depressed Vocabulary",
    "Side-by-side comparison reveals distinct patterns. Depressed group uses mental health terminology (psychiatrist, therapy, depressed) more frequently. Non-depressed group shows more varied conversational vocabulary without clinical terms.",
    f"{img_dir}/02_frequency_comparison.png")

# Slide 12: Correlation Analysis - Method
add_content_slide(prs, "Method: Point-Biserial Correlation", [
    "STATISTICAL APPROACH:",
    "• Point-biserial correlation: measures relationship between",
    "  - Continuous variable (word frequency)",
    "  - Binary variable (depression status: 0 or 1)",
    "",
    "• Calculate correlation coefficient (r) for each word",
    "  - Range: -1 to +1",
    "  - Positive r = word more common in depressed group",
    "  - Negative r = word more common in non-depressed group",
    "  - Magnitude indicates strength of association",
    "",
    "IMPLEMENTATION:",
    "• Use scipy.stats.pointbiserialr function",
    "• Calculate for all 8,474 words in vocabulary",
    "• Select top 30 by absolute correlation strength",
    "• This identifies MOST discriminating words statistically"
])

# Slide 13: Correlation Analysis - Results
add_image_slide(prs, "Results: Words Correlated with Depression",
    "Horizontal bar chart shows strongest correlations. TOP DEPRESSION MARKERS: 'couldnt' (+0.28), 'depressed' (+0.28), 'psychiatrist' (+0.26) - indicating helplessness and clinical awareness. PROTECTIVE WORDS: 'gym', 'siblings', 'travel' (negative r) - indicating active lifestyle and social connections.",
    f"{img_dir}/03_correlations.png")

# Slide 14: TF-IDF Analysis - Method
add_content_slide(prs, "Method: TF-IDF Analysis", [
    "CONCEPT: Term Frequency-Inverse Document Frequency",
    "",
    "TWO COMPONENTS:",
    "1. Term Frequency (TF): How often word appears in a document",
    "2. Inverse Document Frequency (IDF): How rare word is overall",
    "",
    "FORMULA: TF-IDF = TF × IDF",
    "• High TF-IDF = word is frequent in THIS document but rare overall",
    "• Low TF-IDF = word is either rare in document or very common overall",
    "",
    "PURPOSE:",
    "• Identify words that are DISTINCTIVE to specific participants",
    "• Down-weight common words that appear everywhere",
    "• Highlight unique vocabulary that characterizes individuals",
    "",
    "VISUALIZATION: Heatmap showing TF-IDF scores for top words across participants"
])

# Slide 15: TF-IDF - Results
add_image_slide(prs, "Results: TF-IDF Heatmap",
    "Heatmap shows TF-IDF scores (yellow=high, purple=low) for top words across participants. Each row is a participant, each column is a word. Bright yellow spots indicate words that are distinctive to specific individuals, revealing personal vocabulary patterns.",
    f"{img_dir}/06_tfidf_heatmap.png")

# Slide 16: N-gram Analysis - Method
add_content_slide(prs, "Method: N-gram (Bigram) Analysis", [
    "CONCEPT: N-grams are sequences of N consecutive words",
    "• We use BIGRAMS (N=2): two-word phrases",
    "",
    "APPROACH:",
    "• Extract all consecutive word pairs from transcripts",
    "• Count frequency of each bigram",
    "• Compare top bigrams between depressed and non-depressed groups",
    "",
    "WHY BIGRAMS MATTER:",
    "• Single words miss context and meaning",
    "• Phrases reveal patterns: 'dont know' vs 'dont care'",
    "• Capture hesitation markers, common expressions, emotional phrases",
    "",
    "EXPECTED PATTERNS:",
    "• Depressed: more uncertainty phrases ('dont know', 'im not')",
    "• Non-depressed: more active phrases ('like go', 'things like')"
])

# Slide 17: N-gram - Results
add_image_slide(prs, "Results: Bigram Frequency Comparison",
    "Top bigrams for each group shown side-by-side. 'dont know' dominates both groups (uncertainty is universal). Depressed group shows more negative constructions ('cant really', 'feel like'). Non-depressed group uses more casual conversational bigrams.",
    f"{img_dir}/07_ngram_comparison.png")

# Slide 18: Sentiment Analysis - Method
add_content_slide(prs, "Method: VADER Sentiment Analysis", [
    "VADER: Valence Aware Dictionary and sEntiment Reasoner",
    "",
    "HOW IT WORKS:",
    "• Lexicon-based approach (dictionary of ~7,500 words with sentiment scores)",
    "• Analyzes each word and assigns sentiment polarity",
    "• Handles negations, intensifiers, punctuation, capitalization",
    "• Outputs 4 scores for each text:",
    "  - Positive (0-1): proportion of positive sentiment",
    "  - Negative (0-1): proportion of negative sentiment",
    "  - Neutral (0-1): proportion of neutral sentiment",
    "  - Compound (-1 to +1): overall sentiment",
    "",
    "HYPOTHESIS:",
    "• Depressed group will have:",
    "  - Higher negative scores",
    "  - Lower positive scores",
    "  - Lower compound scores (more negative overall)"
])

# Slide 19: Sentiment Analysis - Results
add_content_slide(prs, "Results: Sentiment Score Comparison", [
    "KEY FINDINGS:",
    "• Negative sentiment HIGHER in depressed group",
    "  - Statistical significance: p = 0.0019 (highly significant)",
    "  - Effect size: Cohen's d = -0.559 (medium-large effect)",
    "",
    "• Compound sentiment LOWER in depressed group",
    "  - Statistical significance: p = 0.03 (significant)",
    "  - Depressed group has more negative overall tone",
    "",
    "• Positive sentiment: No significant difference",
    "  - Depression is marked by increased negativity",
    "  - NOT by decreased positivity",
    "",
    "INTERPRETATION:",
    "Language of depressed individuals carries significantly more",
    "negative emotional weight while maintaining neutral tone baseline."
], f"{img_dir}/08_sentiment_comparison.png")

# Slide 20: Word Cloud - Method
add_content_slide(prs, "Method: Word Cloud Visualization", [
    "CONCEPT:",
    "• Visual representation of word frequency",
    "• Larger words = higher frequency in that group",
    "• Color coding for visual appeal and grouping",
    "",
    "APPROACH:",
    "• Generate separate word clouds for each group",
    "• Remove common conversational fillers for clarity",
    "• Size proportional to word frequency",
    "• Provides quick visual summary of vocabulary differences",
    "",
    "ADVANTAGES:",
    "• Intuitive, easy to understand at-a-glance",
    "• Reveals prominent themes and topics",
    "• Complements statistical analysis with visual insight"
])

# Slide 21: Word Cloud - Results
add_two_image_slide(prs, "Results: Word Cloud Comparison",
    "Visual word frequency representation. Larger words = higher frequency in that group. Notice prominent mental health terms in depressed group.",
    f"{img_dir}/09_wordcloud_group0.png",
    f"{img_dir}/09_wordcloud_group1.png",
    "Non-Depressed Group",
    "Depressed Group")

# Slide 22: Key Findings Summary
add_content_slide(prs, "Key Findings: Depression Language Markers", [
    "STRONGEST DEPRESSION INDICATORS (Positive Correlation):",
    "• 'couldnt' (r = +0.278)",
    "  → Expression of helplessness and inability",
    "• 'depressed' (r = +0.277)",
    "  → Direct acknowledgment of condition",
    "• 'psychiatrist' (r = +0.255) / 'therapist' (r = +0.223)",
    "  → Active engagement with mental health treatment",
    "• 'sleep' / 'night' (r = +0.246)",
    "  → Sleep disturbance (key depression symptom)",
    "• 'anxiety' / 'worried' (r = +0.195)",
    "  → Comorbid anxiety symptoms",
    "",
    "These words are STATISTICALLY SIGNIFICANT markers distinguishing",
    "depressed from non-depressed participants (p < 0.01)"
])

# Slide 23: Protective Factors
add_content_slide(prs, "Key Findings: Protective Language Patterns", [
    "PROTECTIVE/NON-DEPRESSED INDICATORS (Negative Correlation):",
    "",
    "SOCIAL CONNECTION:",
    "• 'brother', 'siblings', 'friends' (negative r)",
    "  → Strong family and social relationships",
    "",
    "PHYSICAL ACTIVITY:",
    "• 'gym', 'swim', 'exercise', 'body' (negative r)",
    "  → Active lifestyle and physical health focus",
    "",
    "POSITIVE OUTLOOK:",
    "• 'good', 'pretty', 'forward', 'enjoy' (negative r)",
    "  → Positive framing and forward-looking perspective",
    "",
    "LEISURE & INTERESTS:",
    "• 'travel', 'hobbies', 'music' (negative r)",
    "  → Engagement in enjoyable activities",
    "",
    "These patterns suggest resilience factors and healthy coping mechanisms"
])

# Slide 24: Statistical Validation
add_content_slide(prs, "Statistical Significance & Validation", [
    "POINT-BISERIAL CORRELATION:",
    "• All top markers have p-values < 0.01",
    "• Statistically significant, not due to chance",
    "• Effect sizes range from small to medium",
    "",
    "SENTIMENT ANALYSIS:",
    "• Negative sentiment: p = 0.0019 (highly significant)",
    "• Cohen's d = -0.559 (medium-large effect)",
    "• Compound score: p = 0.03 (significant)",
    "",
    "CLASS IMBALANCE HANDLING:",
    "• Used normalized frequencies (per-person averages)",
    "• Statistical tests account for unequal group sizes",
    "• Results robust despite 30/70 split",
    "",
    "CONCLUSION: Findings are statistically robust and clinically meaningful"
])

# Slide 25: Clinical Implications
add_content_slide(prs, "Clinical & Research Implications", [
    "SCREENING APPLICATIONS:",
    "• Linguistic markers could supplement traditional screening",
    "• Automated analysis of therapy session transcripts",
    "• Early warning system in digital mental health platforms",
    "",
    "TREATMENT INSIGHTS:",
    "• Language patterns may track treatment progress",
    "• Shift from helplessness words to active words = improvement",
    "• Monitor sentiment trends over therapy sessions",
    "",
    "LIMITATIONS:",
    "• Single dataset (AVEC 2017) - needs validation on other datasets",
    "• Cross-sectional (snapshot) - not longitudinal",
    "• Correlation ≠ Causation",
    "• Individual variation is high",
    "",
    "ETHICAL CONSIDERATIONS:",
    "• Privacy of mental health data",
    "• Should augment, not replace, clinical judgment"
])

# Slide 26: Future Directions
add_content_slide(prs, "Future Work & Research Directions", [
    "MACHINE LEARNING CLASSIFICATION:",
    "• Train models to predict depression from language alone",
    "• Test: Logistic Regression, Random Forest, Neural Networks",
    "• Goal: Achieve clinical-grade screening accuracy",
    "",
    "DEEPER LINGUISTIC FEATURES:",
    "• Syntax analysis: sentence structure complexity",
    "• Pronoun usage: self-reference patterns (I, me, my)",
    "• Emotional word categories beyond VADER",
    "• Speech patterns: pauses, hesitations, speech rate",
    "",
    "LONGITUDINAL ANALYSIS:",
    "• Track language changes over treatment course",
    "• Identify markers of improvement vs deterioration",
    "• Personalized monitoring dashboards",
    "",
    "MULTIMODAL INTEGRATION:",
    "• Combine text with audio features (tone, prosody)",
    "• Facial expressions and body language (video)",
    "• Comprehensive depression assessment system"
])

# Slide 27: Conclusion
add_content_slide(prs, "Conclusion", [
    "SUMMARY:",
    "• Successfully identified statistically significant linguistic markers",
    "  of depression in clinical interview transcripts",
    "",
    "• Depressed individuals use:",
    "  - More mental health and treatment terminology",
    "  - More helplessness and inability language",
    "  - Higher negative sentiment overall",
    "",
    "• Non-depressed individuals use:",
    "  - More social connection words",
    "  - More physical activity and leisure terms",
    "  - More positive and forward-looking language",
    "",
    "• Multiple analysis methods (frequency, correlation, TF-IDF,",
    "  N-grams, sentiment) converge on consistent patterns",
    "",
    "                    Thank you! Questions?"
])

# Save
output_path = os.path.join(script_dir, "presentation.pptx")
prs.save(output_path)

print("\n" + "="*60)
print("✓ PowerPoint presentation created successfully!")
print("="*60)
print(f"📄 Output file: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")

# Report missing images if any
if missing_images:
    print(f"\n⚠️  Warning: {len(missing_images)} image(s) not found:")
    for img in missing_images:
        print(f"   - {os.path.basename(img)}")
    print("\n💡 Tip: Run the analysis script to generate all figures.")
else:
    print("✓ All images loaded successfully!")
    
print("="*60 + "\n")
